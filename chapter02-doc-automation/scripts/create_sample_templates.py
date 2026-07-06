#!/usr/bin/env python3
"""테스트용 샘플 회사 템플릿을 생성합니다."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_16x9_template(output_path: str):
    """16:9 회사 템플릿 (빨간 테마)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 테마 색상은 기본 마스터에 포함됨
    # 샘플 슬라이드 1장 추가 (제거 테스트용)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.placeholders[0]
    title.text = "이 슬라이드는 제거됩니다"

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"16:9 템플릿 생성: {output_path}")


def create_4x3_template(output_path: str):
    """4:3 회사 템플릿."""
    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(7.5)

    # 샘플 슬라이드 추가
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.placeholders[0]
    title.text = "4:3 샘플"

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"4:3 템플릿 생성: {output_path}")


if __name__ == "__main__":
    create_16x9_template("../samples/template_16x9.pptx")
    create_4x3_template("../samples/template_4x3.pptx")
