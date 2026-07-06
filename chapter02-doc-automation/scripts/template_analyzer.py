#!/usr/bin/env python3
"""PPT 템플릿 분석 모듈.

회사 PPT 템플릿(.pptx)을 열어 레이아웃, 플레이스홀더, 테마 색상,
슬라이드 크기를 추출합니다.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from lxml import etree


# ──────────────────────────────────────────────
# 테마 색상 → 역할 매핑
# ──────────────────────────────────────────────
THEME_TO_ROLE = {
    "dk1":     "title_text",
    "dk2":     "body_text",
    "lt1":     "card_fill",
    "lt2":     "background",
    "accent1": "accent",
    "accent2": "warning",
    "accent3": "success",
    "accent4": "accent_dark",
    "accent5": "border",
    "accent6": "muted",
}

# 레이아웃 분류 키워드
LAYOUT_KEYWORDS = {
    "cover": ["title slide", "표지", "title only", "제목 슬라이드"],
    "title_content": ["title and content", "제목 및 내용", "title, content"],
    "two_content": ["two content", "2단", "comparison", "비교"],
    "section": ["section header", "구분", "section"],
    "blank": ["blank", "빈 화면", "빈"],
}


# ──────────────────────────────────────────────
# 메인 분석 함수
# ──────────────────────────────────────────────
def analyze_template(template_path: str) -> dict:
    """템플릿 파일을 분석하여 TemplateInfo dict를 반환합니다.

    Returns:
        {
            "slide_width": float (inches),
            "slide_height": float (inches),
            "colors": dict (역할 → hex),
            "font_family": str | None,
            "layouts": dict (용도 → 인덱스),
            "placeholders": dict (인덱스 → 플레이스홀더 목록),
            "existing_slide_count": int,
        }
    """
    prs = Presentation(template_path)

    width, height = _extract_slide_size(prs)
    colors = _extract_theme_colors(prs)
    font_family = _extract_font_family(prs)
    layouts, placeholders = _analyze_layouts(prs)

    return {
        "slide_width": width,
        "slide_height": height,
        "colors": colors,
        "font_family": font_family,
        "layouts": layouts,
        "placeholders": placeholders,
        "existing_slide_count": len(prs.slides),
    }


# ──────────────────────────────────────────────
# 슬라이드 크기 추출
# ──────────────────────────────────────────────
def _extract_slide_size(prs: Presentation) -> tuple[float, float]:
    """슬라이드 크기를 인치로 반환합니다."""
    width_inches = prs.slide_width / Emu(914400)  # 1 inch = 914400 EMU
    height_inches = prs.slide_height / Emu(914400)
    return round(width_inches, 3), round(height_inches, 3)


# ──────────────────────────────────────────────
# 테마 색상 추출
# ──────────────────────────────────────────────
def _extract_theme_colors(prs: Presentation) -> dict:
    """slide_master에서 테마 색상을 추출합니다."""
    colors = {}

    try:
        master = prs.slide_masters[0]
        # XML에서 a:clrScheme 찾기
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = master.element.find(".//a:clrScheme", nsmap)

        if clr_scheme is None:
            return colors

        for theme_key, role in THEME_TO_ROLE.items():
            el = clr_scheme.find(f"a:{theme_key}", nsmap)
            if el is not None:
                # srgbClr (직접 hex) 또는 sysClr (시스템 색상) 처리
                srgb = el.find("a:srgbClr", nsmap)
                if srgb is not None:
                    colors[role] = f"#{srgb.get('val', '000000')}"
                else:
                    sys_clr = el.find("a:sysClr", nsmap)
                    if sys_clr is not None:
                        last_clr = sys_clr.get("lastClr", "000000")
                        colors[role] = f"#{last_clr}"
    except (IndexError, AttributeError):
        pass

    return colors


# ──────────────────────────────────────────────
# 폰트 추출
# ──────────────────────────────────────────────
def _extract_font_family(prs: Presentation) -> str | None:
    """테마의 major/minor 폰트를 추출합니다."""
    try:
        master = prs.slide_masters[0]
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        # minor 폰트 (본문용) 우선
        minor = master.element.find(".//a:fontScheme/a:minorFont/a:latin", nsmap)
        if minor is not None:
            typeface = minor.get("typeface", "")
            if typeface and typeface != "":
                return typeface

        # major 폰트 (제목용)
        major = master.element.find(".//a:fontScheme/a:majorFont/a:latin", nsmap)
        if major is not None:
            typeface = major.get("typeface", "")
            if typeface and typeface != "":
                return typeface
    except (IndexError, AttributeError):
        pass

    return None


# ──────────────────────────────────────────────
# 레이아웃 분석
# ──────────────────────────────────────────────
def _analyze_layouts(prs: Presentation) -> tuple[dict, dict]:
    """모든 slide_layout을 분석하여 용도 매핑과 플레이스홀더 정보를 반환합니다."""
    layouts = {}
    placeholders = {}

    for idx, layout in enumerate(prs.slide_layouts):
        # 레이아웃 분류
        purpose = _classify_layout(layout)
        if purpose and purpose not in layouts:
            layouts[purpose] = idx

        # 플레이스홀더 정보 수집
        ph_list = []
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": _ph_type_name(ph.placeholder_format.type),
                "left": round(ph.left / 914400, 3) if ph.left else 0,
                "top": round(ph.top / 914400, 3) if ph.top else 0,
                "width": round(ph.width / 914400, 3) if ph.width else 0,
                "height": round(ph.height / 914400, 3) if ph.height else 0,
            }
            ph_list.append(ph_info)

        if ph_list:
            placeholders[idx] = ph_list

    # blank 레이아웃이 없으면 마지막 레이아웃을 사용
    if "blank" not in layouts:
        layouts["blank"] = len(prs.slide_layouts) - 1

    return layouts, placeholders


def _classify_layout(layout) -> str | None:
    """레이아웃의 용도를 추론합니다."""
    name = (layout.name or "").lower().strip()

    for purpose, keywords in LAYOUT_KEYWORDS.items():
        for keyword in keywords:
            if keyword in name:
                return purpose

    # 키워드 매칭 실패 시 플레이스홀더 구성으로 판별
    ph_types = set()
    for ph in layout.placeholders:
        ph_types.add(_ph_type_name(ph.placeholder_format.type))

    if ph_types == {"TITLE", "SUBTITLE"} or ph_types == {"CENTER_TITLE", "SUBTITLE"}:
        return "cover"
    if ph_types == {"TITLE", "BODY"} or ph_types == {"TITLE", "OBJECT"}:
        return "title_content"
    if len(ph_types) == 0:
        return "blank"

    return None


def _ph_type_name(ph_type) -> str:
    """플레이스홀더 타입 enum → 문자열 이름."""
    # python-pptx의 PP_PLACEHOLDER enum 값 매핑
    type_map = {
        0: "TITLE",
        1: "BODY",
        2: "CENTER_TITLE",
        3: "SUBTITLE",
        5: "BODY",           # Notes body
        6: "FOOTER",
        7: "DATE",
        8: "OBJECT",
        10: "SLIDE_NUMBER",
        12: "TABLE",
        13: "PICTURE",
        14: "MEDIA_CLIP",
        15: "ORG_CHART",
        16: "CHART",
        18: "PICTURE",       # Bitmap
    }
    try:
        int_val = int(ph_type)
        return type_map.get(int_val, f"OTHER_{int_val}")
    except (ValueError, TypeError):
        return str(ph_type)


# ──────────────────────────────────────────────
# 프리뷰 출력
# ──────────────────────────────────────────────
def print_template_summary(info: dict):
    """템플릿 분석 결과를 읽기 쉽게 출력합니다."""
    print(f"  슬라이드 크기: {info['slide_width']}x{info['slide_height']} 인치", end="")
    ratio = info["slide_width"] / info["slide_height"]
    if abs(ratio - 16 / 9) < 0.1:
        print(" (16:9)")
    elif abs(ratio - 4 / 3) < 0.1:
        print(" (4:3)")
    else:
        print(f" ({ratio:.2f}:1)")

    if info["colors"]:
        print(f"  테마 색상: {len(info['colors'])}개 감지")
        for role, color in list(info["colors"].items())[:5]:
            print(f"    {role}: {color}")
    else:
        print("  테마 색상: 감지 안됨 (기본값 사용)")

    if info["font_family"]:
        print(f"  폰트: {info['font_family']}")

    print(f"  레이아웃: {len(info['layouts'])}개 분류")
    for purpose, idx in info["layouts"].items():
        print(f"    [{idx}] {purpose}")

    if info["existing_slide_count"]:
        print(f"  기존 슬라이드: {info['existing_slide_count']}장 (제거 예정)")


# ──────────────────────────────────────────────
# CLI 실행
# ──────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="PPT 템플릿 분석")
    parser.add_argument("--input", required=True, help="분석할 .pptx 템플릿 경로")
    parser.add_argument("--json", action="store_true", help="JSON 형식으로 출력")
    args = parser.parse_args()

    info = analyze_template(args.input)

    if args.json:
        print(json.dumps(info, ensure_ascii=False, indent=2))
    else:
        print(f"\n[템플릿 분석] {args.input}")
        print("-" * 40)
        print_template_summary(info)


if __name__ == "__main__":
    main()
