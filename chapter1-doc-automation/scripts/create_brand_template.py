#!/usr/bin/env python3
"""브랜드 PPT 템플릿 생성 유틸리티.

회사 로고, 색상, 폰트가 적용된 기본 템플릿 .pptx를 생성합니다.
이 스크립트로 생성된 템플릿을 --template 옵션으로 전달하면
보고서에 브랜드가 자동 적용됩니다.
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_brand_template(
    output_path: str = "templates/brand_template.pptx",
    primary_color: str = "4A90D9",
    secondary_color: str = "2C3E50",
    company_name: str = "My Company",
):
    """브랜드 색상이 적용된 PPT 템플릿을 생성합니다."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary = RGBColor.from_string(primary_color)
    secondary = RGBColor.from_string(secondary_color)

    # 마스터 슬라이드에 브랜드 요소 추가 (빈 슬라이드 1장)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 하단 브랜드 바
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, Inches(7.2), Inches(13.333), Inches(0.3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    # 회사명
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.2), Inches(3), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = company_name
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"브랜드 템플릿 생성: {output_path}")
    print(f"  Primary: #{primary_color}")
    print(f"  Secondary: #{secondary_color}")
    print(f"  Company: {company_name}")


def main():
    parser = argparse.ArgumentParser(description="브랜드 PPT 템플릿 생성")
    parser.add_argument("--output", default="templates/brand_template.pptx",
                        help="출력 파일 경로")
    parser.add_argument("--primary", default="4A90D9",
                        help="주 색상 (hex, 예: 4A90D9)")
    parser.add_argument("--secondary", default="2C3E50",
                        help="보조 색상 (hex, 예: 2C3E50)")
    parser.add_argument("--company", default="My Company",
                        help="회사명")
    args = parser.parse_args()

    create_brand_template(args.output, args.primary, args.secondary, args.company)


if __name__ == "__main__":
    main()
