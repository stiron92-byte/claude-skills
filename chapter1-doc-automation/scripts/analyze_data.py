#!/usr/bin/env python3
"""데이터 분석 및 차트 생성 모듈.

다양한 형식의 파일을 읽고 주요 지표를 계산한 뒤
matplotlib 차트를 생성합니다.

지원 형식:
  - 테이블형: .csv, .xlsx, .xls, .tsv, .json
  - 문서형:   .pdf, .docx, .doc, .pptx, .ppt, .html, .htm, .md, .txt
"""

import argparse
import json
import os
import re
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import pandas as pd

# ──────────────────────────────────────────────
# 한글 폰트 설정
# ──────────────────────────────────────────────
def _setup_korean_font():
    """시스템에서 사용 가능한 한글 폰트를 찾아 설정합니다."""
    korean_fonts = [
        "AppleGothic", "Malgun Gothic", "NanumGothic",
        "NanumBarunGothic", "Noto Sans KR", "Noto Sans CJK KR",
    ]
    available = {f.name for f in fm.fontManager.ttflist}
    for font in korean_fonts:
        if font in available:
            plt.rcParams["font.family"] = font
            break
    plt.rcParams["axes.unicode_minus"] = False

_setup_korean_font()

COLORS = ["#4A90D9", "#E85D75", "#50C878", "#F5A623", "#9B59B6", "#1ABC9C"]

# ──────────────────────────────────────────────
# 파일 형식 분류
# ──────────────────────────────────────────────
TABULAR_EXTENSIONS = {".csv", ".xlsx", ".xls", ".tsv", ".json"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".html", ".htm", ".md", ".txt"}

def detect_file_type(file_path: str) -> str:
    """파일 확장자를 기반으로 'tabular' 또는 'document'를 반환합니다."""
    ext = Path(file_path).suffix.lower()
    if ext in TABULAR_EXTENSIONS:
        return "tabular"
    if ext in DOCUMENT_EXTENSIONS:
        return "document"
    # 알 수 없는 확장자는 텍스트로 시도
    return "document"


# ──────────────────────────────────────────────
# 테이블형 데이터 로더
# ──────────────────────────────────────────────
def load_tabular_data(file_path: str) -> pd.DataFrame:
    """CSV, 엑셀, TSV, JSON 파일을 DataFrame으로 로드합니다."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in (".xlsx", ".xls"):
        return pd.read_excel(path)
    elif ext == ".tsv":
        return pd.read_csv(path, sep="\t")
    elif ext == ".json":
        return pd.read_json(path)
    else:  # .csv 및 기타
        return pd.read_csv(path)


# ──────────────────────────────────────────────
# 문서형 데이터 로더
# ──────────────────────────────────────────────
def load_document_text(file_path: str) -> str:
    """PDF, Word, PPT, HTML, Markdown 등에서 텍스트를 추출합니다."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _load_pdf(path)
    elif ext in (".docx", ".doc"):
        return _load_docx(path)
    elif ext in (".pptx", ".ppt"):
        return _load_pptx(path)
    elif ext in (".html", ".htm"):
        return _load_html(path)
    elif ext == ".md":
        return _load_markdown(path)
    else:  # .txt 및 기타
        return path.read_text(encoding="utf-8")


def _load_pdf(path: Path) -> str:
    """PDF에서 텍스트를 추출합니다."""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        return "\n\n".join(texts)
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            texts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
            return "\n\n".join(texts)
        except ImportError:
            raise ImportError(
                "PDF 파일을 읽으려면 pdfplumber 또는 PyPDF2가 필요합니다.\n"
                "설치: pip install pdfplumber"
            )


def _load_docx(path: Path) -> str:
    """Word 문서에서 텍스트를 추출합니다."""
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # 테이블 내용도 추출
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ImportError(
            "Word 파일을 읽으려면 python-docx가 필요합니다.\n"
            "설치: pip install python-docx"
        )


def _load_pptx(path: Path) -> str:
    """PowerPoint에서 텍스트를 추출합니다."""
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[슬라이드 {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def _load_html(path: Path) -> str:
    """HTML에서 텍스트를 추출합니다."""
    try:
        from bs4 import BeautifulSoup
        html = path.read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "html.parser")
        # 스크립트/스타일 제거
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        # BeautifulSoup 없으면 정규식으로 간단 처리
        html = path.read_text(encoding="utf-8")
        clean = re.sub(r"<[^>]+>", " ", html)
        clean = re.sub(r"\s+", " ", clean)
        return clean.strip()


def _load_markdown(path: Path) -> str:
    """Markdown 파일을 텍스트로 읽습니다."""
    return path.read_text(encoding="utf-8")


# ──────────────────────────────────────────────
# 문서형 분석 (텍스트 기반)
# ──────────────────────────────────────────────
def analyze_document(text: str, file_path: str) -> dict:
    """문서 텍스트를 분석하여 구조화된 결과를 반환합니다."""
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    words = text.split()

    # 숫자 추출 (금액, 수치 등)
    numbers = re.findall(r"[\d,]+\.?\d*", text)
    numbers = [float(n.replace(",", "")) for n in numbers if len(n) > 0]

    # 핵심 문장 추출 (앞부분 위주)
    key_sentences = []
    for line in lines[:50]:
        if len(line) > 10 and not line.startswith(("#", "─", "━", "=", "-")):
            key_sentences.append(line)
        if len(key_sentences) >= 10:
            break

    # 섹션/헤딩 추출
    headings = []
    for line in lines:
        if line.startswith("#") or (len(line) < 80 and line.isupper()):
            headings.append(line.lstrip("#").strip())
    # pptx 슬라이드 헤더
    for line in lines:
        if line.startswith("[슬라이드"):
            headings.append(line)

    result = {
        "input_type": "document",
        "source_file": str(Path(file_path).name),
        "period": {
            "start": "",
            "end": "",
        },
        "summary": {
            "total_characters": len(text),
            "total_lines": len(lines),
            "total_words": len(words),
            "total_sections": len(headings),
            "numbers_found": len(numbers),
        },
        "headings": headings[:20],
        "key_sentences": key_sentences[:10],
        "numbers": sorted(numbers, reverse=True)[:20] if numbers else [],
        "full_text": text[:5000],  # PPT 생성 시 참조용 (최대 5000자)
    }

    return result


# ──────────────────────────────────────────────
# 테이블형 분석 (숫자 기반)
# ──────────────────────────────────────────────
def analyze_sales_data(df: pd.DataFrame) -> dict:
    """매출 데이터를 분석하여 주요 지표를 계산합니다."""
    result = {
        "input_type": "tabular",
        "period": {
            "start": str(df["날짜"].min()) if "날짜" in df.columns else "",
            "end": str(df["날짜"].max()) if "날짜" in df.columns else "",
        },
        "summary": {},
        "daily": {},
        "category": {},
        "top_items": [],
    }

    # 숫자 컬럼 자동 감지
    numeric_cols = df.select_dtypes(include="number").columns.tolist()

    # 날짜 + 이름 + 매출 + 판매량 컬럼이 있는 경우 (매출 분석)
    if "매출" in df.columns and "판매량" in df.columns:
        total_revenue = df["매출"].sum()
        total_quantity = df["판매량"].sum()

        date_col = "날짜" if "날짜" in df.columns else None
        name_col = "상품명" if "상품명" in df.columns else None

        avg_daily = 0
        if date_col:
            avg_daily = df.groupby(date_col)["매출"].sum().mean()

        result["summary"] = {
            "total_revenue": int(total_revenue),
            "total_quantity": int(total_quantity),
            "avg_daily_revenue": int(avg_daily),
            "unique_products": int(df[name_col].nunique()) if name_col else 0,
            "total_days": int(df[date_col].nunique()) if date_col else 0,
        }

        if date_col:
            daily = df.groupby(date_col).agg({"매출": "sum", "판매량": "sum"}).reset_index()
            result["daily"] = {
                "dates": daily[date_col].tolist(),
                "revenue": daily["매출"].tolist(),
                "quantity": daily["판매량"].tolist(),
            }
            # 전주 대비 변화율
            mid = len(daily) // 2
            if mid > 0:
                first_half = daily.iloc[:mid]["매출"].sum()
                second_half = daily.iloc[mid:]["매출"].sum()
                if first_half > 0:
                    change_rate = ((second_half - first_half) / first_half) * 100
                else:
                    change_rate = 0.0
                result["summary"]["revenue_change_rate"] = round(change_rate, 1)

        if "카테고리" in df.columns:
            cat = df.groupby("카테고리").agg({"매출": "sum", "판매량": "sum"}).reset_index()
            cat = cat.sort_values("매출", ascending=False)
            result["category"] = {
                "names": cat["카테고리"].tolist(),
                "revenue": cat["매출"].tolist(),
                "quantity": cat["판매량"].tolist(),
            }

        if name_col:
            top = df.groupby(name_col).agg({"매출": "sum", "판매량": "sum"}).reset_index()
            top = top.sort_values("매출", ascending=False).head(5)
            result["top_items"] = [
                {"name": row[name_col], "revenue": int(row["매출"]), "quantity": int(row["판매량"])}
                for _, row in top.iterrows()
            ]

    else:
        # 범용 테이블 분석 (매출/판매량 컬럼이 없는 경우)
        result["summary"] = {
            "total_rows": len(df),
            "total_columns": len(df.columns),
            "columns": df.columns.tolist(),
            "numeric_columns": numeric_cols,
        }
        # 숫자 컬럼별 통계
        for col in numeric_cols[:5]:
            result["summary"][f"{col}_total"] = float(df[col].sum())
            result["summary"][f"{col}_mean"] = float(df[col].mean())
            result["summary"][f"{col}_max"] = float(df[col].max())

    return result


# ──────────────────────────────────────────────
# 차트 생성
# ──────────────────────────────────────────────
def generate_charts(analysis: dict, output_dir: str) -> list[str]:
    """분석 결과를 기반으로 차트를 생성합니다."""
    charts_dir = Path(output_dir) / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)
    chart_paths = []

    # 문서형은 차트 생성 스킵
    if analysis.get("input_type") == "document":
        return chart_paths

    # 일별 매출 추이 라인 차트
    if analysis.get("daily", {}).get("dates"):
        fig, ax = plt.subplots(figsize=(10, 5))
        dates = analysis["daily"]["dates"]
        revenue = analysis["daily"]["revenue"]
        ax.plot(dates, revenue, color=COLORS[0], linewidth=2, marker="o", markersize=5)
        ax.fill_between(dates, revenue, alpha=0.1, color=COLORS[0])
        ax.set_title("일별 매출 추이", fontsize=16, fontweight="bold", pad=15)
        ax.set_xlabel("날짜", fontsize=12)
        ax.set_ylabel("매출 (원)", fontsize=12)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:,.0f}"))
        plt.xticks(rotation=45, ha="right")
        plt.tight_layout()
        path = str(charts_dir / "daily_revenue.png")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        chart_paths.append(path)

    # 카테고리별 매출 파이 차트
    if analysis.get("category", {}).get("names"):
        fig, ax = plt.subplots(figsize=(8, 8))
        names = analysis["category"]["names"]
        values = analysis["category"]["revenue"]
        colors = COLORS[: len(names)]
        wedges, texts, autotexts = ax.pie(
            values, labels=names, colors=colors, autopct="%1.1f%%",
            startangle=90, pctdistance=0.85,
        )
        for text in autotexts:
            text.set_fontsize(11)
            text.set_fontweight("bold")
        ax.set_title("카테고리별 매출 비중", fontsize=16, fontweight="bold", pad=15)
        plt.tight_layout()
        path = str(charts_dir / "category_pie.png")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        chart_paths.append(path)

    # TOP 5 상품 막대 차트
    if analysis.get("top_items"):
        fig, ax = plt.subplots(figsize=(10, 5))
        items = analysis["top_items"]
        names = [item["name"] for item in items]
        revenues = [item["revenue"] for item in items]
        bars = ax.barh(names[::-1], revenues[::-1], color=COLORS[:len(names)])
        for bar in bars:
            width = bar.get_width()
            ax.text(width + max(revenues) * 0.02, bar.get_y() + bar.get_height() / 2,
                    f"{width:,.0f}원", va="center", fontsize=10)
        ax.set_title("TOP 5 상품 매출", fontsize=16, fontweight="bold", pad=15)
        ax.set_xlabel("매출 (원)", fontsize=12)
        ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:,.0f}"))
        plt.tight_layout()
        path = str(charts_dir / "top_products.png")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        chart_paths.append(path)

    return chart_paths


# ──────────────────────────────────────────────
# 메인 실행
# ──────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="데이터 분석 및 차트 생성")
    parser.add_argument("--input", required=True, help="입력 파일 경로")
    parser.add_argument("--output", default="output", help="출력 디렉토리")
    args = parser.parse_args()

    file_type = detect_file_type(args.input)
    ext = Path(args.input).suffix.lower()
    print(f"[1/3] 파일 로드: {args.input} ({ext}, {file_type})")

    if file_type == "tabular":
        df = load_tabular_data(args.input)
        print(f"      → {len(df)}행 x {len(df.columns)}열 로드 완료")
        print("[2/3] 데이터 분석 중...")
        analysis = analyze_sales_data(df)
        summary = analysis["summary"]
        if "total_revenue" in summary:
            print(f"      → 총 매출: {summary['total_revenue']:,}원")
            print(f"      → 총 판매량: {summary['total_quantity']:,}개")
            if "revenue_change_rate" in summary:
                rate = summary["revenue_change_rate"]
                arrow = "↑" if rate >= 0 else "↓"
                print(f"      → 전주 대비: {arrow} {abs(rate)}%")
        else:
            print(f"      → {summary.get('total_rows', 0)}행 x {summary.get('total_columns', 0)}열")
    else:
        text = load_document_text(args.input)
        print(f"      → {len(text):,}자 텍스트 추출 완료")
        print("[2/3] 문서 분석 중...")
        analysis = analyze_document(text, args.input)
        summary = analysis["summary"]
        print(f"      → {summary['total_lines']}줄, {summary['total_words']}단어")
        print(f"      → 섹션 {summary['total_sections']}개, 숫자 {summary['numbers_found']}개 감지")

    print("[3/3] 차트 생성 중...")
    chart_paths = generate_charts(analysis, args.output)
    if chart_paths:
        print(f"      → {len(chart_paths)}개 차트 생성 완료")
    else:
        print(f"      → 문서형 입력 — 차트 생성 건너뜀")

    # 분석 결과 저장
    output_path = Path(args.output)
    output_path.mkdir(parents=True, exist_ok=True)
    analysis_file = str(output_path / "analysis.json")
    analysis["chart_paths"] = chart_paths
    # full_text는 JSON 저장 시 별도 파일로 분리
    full_text = analysis.pop("full_text", None)
    with open(analysis_file, "w", encoding="utf-8") as f:
        json.dump(analysis, f, ensure_ascii=False, indent=2)
    if full_text:
        text_file = str(output_path / "extracted_text.txt")
        Path(text_file).write_text(full_text, encoding="utf-8")
        analysis["full_text_path"] = text_file
        print(f"      추출 텍스트: {text_file}")
    print(f"\n분석 결과 저장: {analysis_file}")

    return analysis


if __name__ == "__main__":
    main()
