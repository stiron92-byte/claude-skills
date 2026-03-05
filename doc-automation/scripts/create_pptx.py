#!/usr/bin/env python3
"""PPT 생성 모듈.

분석 결과와 차트를 기반으로 python-pptx를 사용하여
브랜드 템플릿이 적용된 .pptx 파일을 생성합니다.
"""

import argparse
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 브랜드 색상
BRAND_PRIMARY = RGBColor(0x4A, 0x90, 0xD9)    # 파란색
BRAND_SECONDARY = RGBColor(0x2C, 0x3E, 0x50)  # 진한 남색
BRAND_ACCENT = RGBColor(0xE8, 0x5D, 0x75)     # 핑크/레드
BRAND_SUCCESS = RGBColor(0x50, 0xC8, 0x78)     # 초록
BRAND_BG = RGBColor(0xF8, 0xF9, 0xFA)          # 밝은 회색
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x6C, 0x75, 0x7D)


def _add_bg_shape(slide, color=BRAND_BG):
    """슬라이드 배경 사각형을 추가합니다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_header_bar(slide, color=BRAND_PRIMARY):
    """상단 헤더 바를 추가합니다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(1.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """텍스트 박스를 추가합니다."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_kpi_card(slide, left, top, width, label, value, change=None):
    """KPI 카드(박스형 지표)를 추가합니다."""
    # 카드 배경
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(1.8),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
    shape.line.width = Pt(1)

    # 라벨
    _add_text_box(slide, left + 0.2, top + 0.2, width - 0.4, 0.4,
                  label, font_size=11, color=GRAY)

    # 값
    _add_text_box(slide, left + 0.2, top + 0.6, width - 0.4, 0.6,
                  value, font_size=24, bold=True, color=BRAND_SECONDARY)

    # 변화율
    if change is not None:
        arrow = "▲" if change >= 0 else "▼"
        change_color = BRAND_SUCCESS if change >= 0 else BRAND_ACCENT
        _add_text_box(slide, left + 0.2, top + 1.2, width - 0.4, 0.4,
                      f"{arrow} {abs(change)}% 전주 대비", font_size=10, color=change_color)


def create_cover_slide(prs, analysis):
    """표지 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg_shape(slide, BRAND_PRIMARY)

    # 제목
    _add_text_box(slide, 1.5, 1.5, 10, 1.2,
                  "주간 업무 보고서", font_size=44, bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)

    # 기간
    period = analysis.get("period", {})
    period_text = f"{period.get('start', '')} ~ {period.get('end', '')}"
    _add_text_box(slide, 1.5, 3.0, 10, 0.6,
                  period_text, font_size=20, color=WHITE,
                  alignment=PP_ALIGN.CENTER)

    # 구분선
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(3.8), Inches(3.333), Inches(0.03),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # 작성일/작성자
    _add_text_box(slide, 1.5, 4.2, 10, 0.5,
                  f"작성일: {datetime.now().strftime('%Y-%m-%d')}",
                  font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)


def create_summary_slide(prs, analysis):
    """핵심 KPI 요약 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_shape(slide)
    _add_header_bar(slide)

    # 헤더 텍스트
    _add_text_box(slide, 0.5, 0.15, 12, 0.7,
                  "핵심 지표 요약", font_size=28, bold=True, color=WHITE)

    summary = analysis.get("summary", {})
    change_rate = summary.get("revenue_change_rate")

    # KPI 카드들
    _add_kpi_card(slide, 0.8, 1.5, 3.8,
                  "총 매출", f"{summary.get('total_revenue', 0):,}원",
                  change=change_rate)
    _add_kpi_card(slide, 5.0, 1.5, 3.8,
                  "총 판매량", f"{summary.get('total_quantity', 0):,}개")
    _add_kpi_card(slide, 9.2, 1.5, 3.8,
                  "일 평균 매출", f"{summary.get('avg_daily_revenue', 0):,}원")

    # TOP 5 테이블
    _add_text_box(slide, 0.8, 3.8, 12, 0.5,
                  "TOP 5 상품", font_size=18, bold=True, color=BRAND_SECONDARY)

    top_items = analysis.get("top_items", [])
    y = 4.4
    for i, item in enumerate(top_items[:5], 1):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y), Inches(11.7), Inches(0.45),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE if i % 2 == 1 else BRAND_BG
        bg.line.fill.background()

        _add_text_box(slide, 1.0, y + 0.05, 0.5, 0.35,
                      f"{i}", font_size=12, bold=True, color=BRAND_PRIMARY)
        _add_text_box(slide, 1.6, y + 0.05, 5, 0.35,
                      item["name"], font_size=12, color=DARK)
        _add_text_box(slide, 8.0, y + 0.05, 2.5, 0.35,
                      f"{item['revenue']:,}원", font_size=12, bold=True,
                      color=BRAND_SECONDARY, alignment=PP_ALIGN.RIGHT)
        _add_text_box(slide, 10.8, y + 0.05, 1.5, 0.35,
                      f"{item['quantity']:,}개", font_size=12,
                      color=GRAY, alignment=PP_ALIGN.RIGHT)
        y += 0.5


def create_chart_slide(prs, analysis):
    """차트 포함 상세 분석 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_shape(slide)
    _add_header_bar(slide)

    _add_text_box(slide, 0.5, 0.15, 12, 0.7,
                  "상세 분석", font_size=28, bold=True, color=WHITE)

    chart_paths = analysis.get("chart_paths", [])

    if len(chart_paths) >= 1 and Path(chart_paths[0]).exists():
        slide.shapes.add_picture(
            chart_paths[0],
            Inches(0.5), Inches(1.3), Inches(7.5), Inches(3.5),
        )

    if len(chart_paths) >= 2 and Path(chart_paths[1]).exists():
        slide.shapes.add_picture(
            chart_paths[1],
            Inches(8.2), Inches(1.3), Inches(4.5), Inches(4.5),
        )

    if len(chart_paths) >= 3 and Path(chart_paths[2]).exists():
        slide.shapes.add_picture(
            chart_paths[2],
            Inches(0.5), Inches(5.0), Inches(7.5), Inches(2.3),
        )


def create_issues_slide(prs, analysis):
    """주요 이슈 및 액션 아이템 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_shape(slide)
    _add_header_bar(slide)

    _add_text_box(slide, 0.5, 0.15, 12, 0.7,
                  "주요 이슈 & 액션 아이템", font_size=28, bold=True, color=WHITE)

    summary = analysis.get("summary", {})
    change_rate = summary.get("revenue_change_rate", 0)

    # 이슈 섹션
    _add_text_box(slide, 0.8, 1.3, 5.5, 0.5,
                  "주요 이슈", font_size=20, bold=True, color=BRAND_SECONDARY)

    issues = []
    if change_rate < 0:
        issues.append(f"매출 전주 대비 {abs(change_rate)}% 감소 — 원인 분석 필요")
    elif change_rate > 10:
        issues.append(f"매출 전주 대비 {change_rate}% 성장 — 성공 요인 분석")

    top_items = analysis.get("top_items", [])
    if top_items:
        issues.append(f"'{top_items[0]['name']}' 매출 1위 — 재고 관리 주의")
    issues.append("데이터 기반 의사결정 체계 개선 필요")

    y = 1.9
    for issue in issues:
        # 아이콘 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), Inches(y + 0.05), Inches(0.25), Inches(0.25),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BRAND_ACCENT
        circle.line.fill.background()
        _add_text_box(slide, 1.5, y, 5, 0.4, issue, font_size=13, color=DARK)
        y += 0.55

    # 액션 아이템 섹션
    _add_text_box(slide, 7.0, 1.3, 5.5, 0.5,
                  "액션 아이템", font_size=20, bold=True, color=BRAND_SECONDARY)

    actions = [
        "매출 변동 원인 심층 분석",
        "TOP 상품 프로모션 전략 수립",
        "재고 최적화 회의 진행",
        "다음주 목표 KPI 설정",
    ]
    y = 1.9
    for action in actions:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.2), Inches(y + 0.05), Inches(0.25), Inches(0.25),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BRAND_SUCCESS
        circle.line.fill.background()
        _add_text_box(slide, 7.7, y, 5, 0.4, action, font_size=13, color=DARK)
        y += 0.55


def create_plan_slide(prs, analysis):
    """다음 주 계획 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_shape(slide)
    _add_header_bar(slide, BRAND_SECONDARY)

    _add_text_box(slide, 0.5, 0.15, 12, 0.7,
                  "다음 주 계획", font_size=28, bold=True, color=WHITE)

    summary = analysis.get("summary", {})
    current_revenue = summary.get("total_revenue", 0)
    target = int(current_revenue * 1.05)  # 5% 성장 목표

    # 목표 섹션
    _add_text_box(slide, 0.8, 1.3, 12, 0.5,
                  "주간 목표", font_size=20, bold=True, color=BRAND_SECONDARY)

    _add_kpi_card(slide, 0.8, 1.9, 3.8,
                  "매출 목표", f"{target:,}원")
    _add_kpi_card(slide, 5.0, 1.9, 3.8,
                  "성장률 목표", "5%")
    _add_kpi_card(slide, 9.2, 1.9, 3.8,
                  "신규 고객 목표", "+10명")

    # 일정
    _add_text_box(slide, 0.8, 4.2, 12, 0.5,
                  "주요 일정", font_size=20, bold=True, color=BRAND_SECONDARY)

    schedule = [
        ("월", "주간 킥오프 미팅 & KPI 리뷰"),
        ("화~수", "프로모션 캠페인 기획 및 실행"),
        ("목", "중간 점검 및 전략 조정"),
        ("금", "주간 마감 및 보고서 작성"),
    ]

    y = 4.8
    for day, task in schedule:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y), Inches(11.7), Inches(0.45),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        bg.line.width = Pt(1)

        _add_text_box(slide, 1.0, y + 0.05, 1.5, 0.35,
                      day, font_size=13, bold=True, color=BRAND_PRIMARY)
        _add_text_box(slide, 2.8, y + 0.05, 9, 0.35,
                      task, font_size=13, color=DARK)
        y += 0.55


def create_presentation(analysis: dict, output_path: str, template_path: str = None):
    """분석 결과를 기반으로 PPT를 생성합니다."""
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    print("  [슬라이드 1/5] 표지 생성...")
    create_cover_slide(prs, analysis)

    print("  [슬라이드 2/5] 핵심 KPI 요약...")
    create_summary_slide(prs, analysis)

    print("  [슬라이드 3/5] 상세 분석 차트...")
    create_chart_slide(prs, analysis)

    print("  [슬라이드 4/5] 주요 이슈 & 액션...")
    create_issues_slide(prs, analysis)

    print("  [슬라이드 5/5] 다음 주 계획...")
    create_plan_slide(prs, analysis)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"\n  PPT 저장 완료: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="PPT 보고서 생성")
    parser.add_argument("--analysis", required=True, help="analysis.json 경로")
    parser.add_argument("--charts", default="output/charts", help="차트 디렉토리")
    parser.add_argument("--template", default=None, help="브랜드 템플릿 .pptx 경로")
    parser.add_argument("--output", default="output/weekly_report.pptx", help="출력 파일 경로")
    args = parser.parse_args()

    with open(args.analysis, "r", encoding="utf-8") as f:
        analysis = json.load(f)

    print("[PPT 생성 시작]")
    create_presentation(analysis, args.output, args.template)


if __name__ == "__main__":
    main()
