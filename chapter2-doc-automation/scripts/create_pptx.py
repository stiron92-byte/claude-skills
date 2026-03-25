#!/usr/bin/env python3
"""PPT 생성 모듈.

분석 결과와 차트를 기반으로 python-pptx를 사용하여
브랜드 템플릿이 적용된 .pptx 파일을 생성합니다.

StyleConfig를 통해 색상, 폰트, 좌표를 관리하며,
회사 템플릿의 플레이스홀더가 있으면 우선 사용합니다.
"""

import argparse
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from style_config import StyleConfig, default_style


# ──────────────────────────────────────────────
# 헬퍼 함수
# ──────────────────────────────────────────────
def _add_bg_shape(slide, style: StyleConfig, color_role: str = "background"):
    """슬라이드 배경 사각형을 추가합니다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(style.slide_width), Inches(style.slide_height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.rgb(color_role)
    shape.line.fill.background()


def _add_header_bar(slide, style: StyleConfig, color_role: str = "accent"):
    """상단 헤더 바를 추가합니다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(style.slide_width), Inches(style.y(0.133)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.rgb(color_role)
    shape.line.fill.background()


def _add_text_box(slide, left, top, width, height, text,
                  font_size=14, bold=False, color=None,
                  alignment=PP_ALIGN.LEFT, font_family=None):
    """텍스트 박스를 추가합니다."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    if font_family:
        p.font.name = font_family
    return txBox


def _add_kpi_card(slide, left, top, width, label, value, style: StyleConfig, change=None):
    """KPI 카드(박스형 지표)를 추가합니다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(style.y(0.24)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.rgb("card_fill")
    shape.line.color.rgb = style.rgb("border")
    shape.line.width = Pt(1)

    _add_text_box(slide, left + 0.2, top + 0.2, width - 0.4, 0.4,
                  label, font_size=style.font("caption")["size_pt"],
                  color=style.rgb("muted"), font_family=style.font_family)

    _add_text_box(slide, left + 0.2, top + 0.6, width - 0.4, 0.6,
                  value, font_size=24, bold=True,
                  color=style.rgb("accent_dark"), font_family=style.font_family)

    if change is not None:
        arrow = "▲" if change >= 0 else "▼"
        change_color = style.rgb("success") if change >= 0 else style.rgb("warning")
        _add_text_box(slide, left + 0.2, top + 1.2, width - 0.4, 0.4,
                      f"{arrow} {abs(change)}% 전주 대비",
                      font_size=style.font("small")["size_pt"],
                      color=change_color, font_family=style.font_family)


def _find_placeholder(slide, ph_type: str):
    """슬라이드에서 특정 타입의 플레이스홀더를 찾습니다. 없으면 None."""
    type_map = {
        "TITLE": (0, 2),       # TITLE, CENTER_TITLE
        "SUBTITLE": (3,),
        "BODY": (1, 8),        # BODY, OBJECT
        "PICTURE": (13, 18),
    }
    target_ids = type_map.get(ph_type, ())
    for ph in slide.placeholders:
        try:
            if int(ph.placeholder_format.type) in target_ids:
                return ph
        except (ValueError, TypeError):
            pass
    return None


def _apply_placeholder_text(ph, text, style: StyleConfig, font_level: str, color_role: str):
    """플레이스홀더에 텍스트와 스타일을 적용합니다."""
    ph.text = text
    font_cfg = style.font(font_level)
    for paragraph in ph.text_frame.paragraphs:
        paragraph.font.size = Pt(font_cfg["size_pt"])
        paragraph.font.bold = font_cfg["bold"]
        paragraph.font.color.rgb = style.rgb(color_role)
        if style.font_family:
            paragraph.font.name = style.font_family


def _get_layout(prs, style: StyleConfig, purpose: str) -> int:
    """용도에 맞는 레이아웃 인덱스를 반환합니다."""
    idx = style.get_layout_index(purpose)
    # 인덱스가 유효한지 확인
    if idx < len(prs.slide_layouts):
        return idx
    return min(6, len(prs.slide_layouts) - 1)


def _remove_existing_slides(prs):
    """템플릿의 기존 슬라이드를 모두 제거합니다. 마스터/레이아웃은 유지."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ──────────────────────────────────────────────
# 슬라이드 생성 함수
# ──────────────────────────────────────────────
def create_cover_slide(prs, analysis, style: StyleConfig):
    """표지 슬라이드를 생성합니다."""
    layout_idx = _get_layout(prs, style, "cover")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # 플레이스홀더 시도
    title_ph = _find_placeholder(slide, "TITLE")
    subtitle_ph = _find_placeholder(slide, "SUBTITLE")

    if title_ph:
        _apply_placeholder_text(title_ph, "주간 업무 보고서", style, "display", "card_fill")
    else:
        _add_bg_shape(slide, style, "accent")
        _add_text_box(slide, style.x(0.112), style.y(0.2), style.x(0.75), style.y(0.16),
                      "주간 업무 보고서",
                      font_size=style.font("display")["size_pt"],
                      bold=True, color=style.rgb("card_fill"),
                      alignment=PP_ALIGN.CENTER, font_family=style.font_family)

    period = analysis.get("period", {})
    period_text = f"{period.get('start', '')} ~ {period.get('end', '')}"

    if subtitle_ph:
        _apply_placeholder_text(subtitle_ph, period_text, style, "subhead", "card_fill")
    elif not title_ph:
        _add_text_box(slide, style.x(0.112), style.y(0.4), style.x(0.75), style.y(0.08),
                      period_text,
                      font_size=style.font("subhead")["size_pt"],
                      color=style.rgb("card_fill"),
                      alignment=PP_ALIGN.CENTER, font_family=style.font_family)

        # 구분선
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(style.x(0.375)), Inches(style.y(0.507)),
            Inches(style.x(0.25)), Inches(style.y(0.004)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style.rgb("card_fill")
        shape.line.fill.background()

        # 작성일
        _add_text_box(slide, style.x(0.112), style.y(0.56), style.x(0.75), style.y(0.067),
                      f"작성일: {datetime.now().strftime('%Y-%m-%d')}",
                      font_size=style.font("body")["size_pt"],
                      color=style.rgb("card_fill"),
                      alignment=PP_ALIGN.CENTER, font_family=style.font_family)


def create_summary_slide(prs, analysis, style: StyleConfig):
    """핵심 KPI 요약 슬라이드를 생성합니다."""
    layout_idx = _get_layout(prs, style, "blank")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title_ph = _find_placeholder(slide, "TITLE")
    if title_ph:
        _apply_placeholder_text(title_ph, "핵심 지표 요약", style, "heading", "card_fill")
    else:
        _add_bg_shape(slide, style)
        _add_header_bar(slide, style)
        _add_text_box(slide, style.x(0.0375), style.y(0.02), style.x(0.9), style.y(0.093),
                      "핵심 지표 요약",
                      font_size=style.font("heading")["size_pt"],
                      bold=True, color=style.rgb("card_fill"),
                      font_family=style.font_family)

    summary = analysis.get("summary", {})
    change_rate = summary.get("revenue_change_rate")

    # KPI 카드 — 비율 좌표
    card_y = style.y(0.2)
    card_w = style.x(0.285)
    _add_kpi_card(slide, style.x(0.06), card_y, card_w,
                  "총 매출", f"{summary.get('total_revenue', 0):,}원",
                  style, change=change_rate)
    _add_kpi_card(slide, style.x(0.375), card_y, card_w,
                  "총 판매량", f"{summary.get('total_quantity', 0):,}개", style)
    _add_kpi_card(slide, style.x(0.69), card_y, card_w,
                  "일 평균 매출", f"{summary.get('avg_daily_revenue', 0):,}원", style)

    # TOP 5
    _add_text_box(slide, style.x(0.06), style.y(0.507), style.x(0.9), style.y(0.067),
                  "TOP 5 상품",
                  font_size=style.font("subhead")["size_pt"] - 2,
                  bold=True, color=style.rgb("accent_dark"),
                  font_family=style.font_family)

    top_items = analysis.get("top_items", [])
    row_y = style.y(0.587)
    row_h = style.y(0.06)
    for i, item in enumerate(top_items[:5], 1):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(style.x(0.06)), Inches(row_y),
            Inches(style.x(0.877)), Inches(row_h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = style.rgb("card_fill") if i % 2 == 1 else style.rgb("background")
        bg.line.fill.background()

        _add_text_box(slide, style.x(0.075), row_y + style.y(0.007), style.x(0.0375), row_h * 0.8,
                      f"{i}", font_size=12, bold=True, color=style.rgb("accent"),
                      font_family=style.font_family)
        _add_text_box(slide, style.x(0.12), row_y + style.y(0.007), style.x(0.375), row_h * 0.8,
                      item["name"], font_size=12, color=style.rgb("body_text"),
                      font_family=style.font_family)
        _add_text_box(slide, style.x(0.6), row_y + style.y(0.007), style.x(0.187), row_h * 0.8,
                      f"{item['revenue']:,}원", font_size=12, bold=True,
                      color=style.rgb("accent_dark"), alignment=PP_ALIGN.RIGHT,
                      font_family=style.font_family)
        _add_text_box(slide, style.x(0.81), row_y + style.y(0.007), style.x(0.112), row_h * 0.8,
                      f"{item['quantity']:,}개", font_size=12,
                      color=style.rgb("muted"), alignment=PP_ALIGN.RIGHT,
                      font_family=style.font_family)
        row_y += row_h + style.y(0.007)


def create_chart_slide(prs, analysis, style: StyleConfig):
    """차트 포함 상세 분석 슬라이드를 생성합니다."""
    layout_idx = _get_layout(prs, style, "blank")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title_ph = _find_placeholder(slide, "TITLE")
    if title_ph:
        _apply_placeholder_text(title_ph, "상세 분석", style, "heading", "card_fill")
    else:
        _add_bg_shape(slide, style)
        _add_header_bar(slide, style)
        _add_text_box(slide, style.x(0.0375), style.y(0.02), style.x(0.9), style.y(0.093),
                      "상세 분석",
                      font_size=style.font("heading")["size_pt"],
                      bold=True, color=style.rgb("card_fill"),
                      font_family=style.font_family)

    chart_paths = analysis.get("chart_paths", [])

    if len(chart_paths) >= 1 and Path(chart_paths[0]).exists():
        slide.shapes.add_picture(
            chart_paths[0],
            Inches(style.x(0.0375)), Inches(style.y(0.173)),
            Inches(style.x(0.5625)), Inches(style.y(0.467)),
        )

    if len(chart_paths) >= 2 and Path(chart_paths[1]).exists():
        slide.shapes.add_picture(
            chart_paths[1],
            Inches(style.x(0.615)), Inches(style.y(0.173)),
            Inches(style.x(0.3375)), Inches(style.y(0.6)),
        )

    if len(chart_paths) >= 3 and Path(chart_paths[2]).exists():
        slide.shapes.add_picture(
            chart_paths[2],
            Inches(style.x(0.0375)), Inches(style.y(0.667)),
            Inches(style.x(0.5625)), Inches(style.y(0.307)),
        )


def create_issues_slide(prs, analysis, style: StyleConfig):
    """주요 이슈 및 액션 아이템 슬라이드를 생성합니다."""
    layout_idx = _get_layout(prs, style, "blank")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title_ph = _find_placeholder(slide, "TITLE")
    if title_ph:
        _apply_placeholder_text(title_ph, "주요 이슈 & 액션 아이템", style, "heading", "card_fill")
    else:
        _add_bg_shape(slide, style)
        _add_header_bar(slide, style)
        _add_text_box(slide, style.x(0.0375), style.y(0.02), style.x(0.9), style.y(0.093),
                      "주요 이슈 & 액션 아이템",
                      font_size=style.font("heading")["size_pt"],
                      bold=True, color=style.rgb("card_fill"),
                      font_family=style.font_family)

    summary = analysis.get("summary", {})
    change_rate = summary.get("revenue_change_rate", 0)

    # 이슈 섹션
    _add_text_box(slide, style.x(0.06), style.y(0.173), style.x(0.41), style.y(0.067),
                  "주요 이슈",
                  font_size=style.font("subhead")["size_pt"],
                  bold=True, color=style.rgb("accent_dark"),
                  font_family=style.font_family)

    issues = []
    if change_rate < 0:
        issues.append(f"매출 전주 대비 {abs(change_rate)}% 감소 — 원인 분석 필요")
    elif change_rate > 10:
        issues.append(f"매출 전주 대비 {change_rate}% 성장 — 성공 요인 분석")
    top_items = analysis.get("top_items", [])
    if top_items:
        issues.append(f"'{top_items[0]['name']}' 매출 1위 — 재고 관리 주의")
    issues.append("데이터 기반 의사결정 체계 개선 필요")

    bullet_y = style.y(0.253)
    for issue in issues:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(style.x(0.075)), Inches(bullet_y + style.y(0.007)),
            Inches(style.x(0.019)), Inches(style.x(0.019)),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = style.rgb("warning")
        circle.line.fill.background()
        _add_text_box(slide, style.x(0.112), bullet_y, style.x(0.375), style.y(0.053),
                      issue, font_size=13, color=style.rgb("body_text"),
                      font_family=style.font_family)
        bullet_y += style.y(0.073)

    # 액션 아이템 섹션
    _add_text_box(slide, style.x(0.525), style.y(0.173), style.x(0.41), style.y(0.067),
                  "액션 아이템",
                  font_size=style.font("subhead")["size_pt"],
                  bold=True, color=style.rgb("accent_dark"),
                  font_family=style.font_family)

    actions = [
        "매출 변동 원인 심층 분석",
        "TOP 상품 프로모션 전략 수립",
        "재고 최적화 회의 진행",
        "다음주 목표 KPI 설정",
    ]
    bullet_y = style.y(0.253)
    for action in actions:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(style.x(0.54)), Inches(bullet_y + style.y(0.007)),
            Inches(style.x(0.019)), Inches(style.x(0.019)),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = style.rgb("success")
        circle.line.fill.background()
        _add_text_box(slide, style.x(0.577), bullet_y, style.x(0.375), style.y(0.053),
                      action, font_size=13, color=style.rgb("body_text"),
                      font_family=style.font_family)
        bullet_y += style.y(0.073)


def create_plan_slide(prs, analysis, style: StyleConfig):
    """다음 주 계획 슬라이드를 생성합니다."""
    layout_idx = _get_layout(prs, style, "blank")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title_ph = _find_placeholder(slide, "TITLE")
    if title_ph:
        _apply_placeholder_text(title_ph, "다음 주 계획", style, "heading", "card_fill")
    else:
        _add_bg_shape(slide, style)
        _add_header_bar(slide, style, "accent_dark")
        _add_text_box(slide, style.x(0.0375), style.y(0.02), style.x(0.9), style.y(0.093),
                      "다음 주 계획",
                      font_size=style.font("heading")["size_pt"],
                      bold=True, color=style.rgb("card_fill"),
                      font_family=style.font_family)

    summary = analysis.get("summary", {})
    current_revenue = summary.get("total_revenue", 0)
    target = int(current_revenue * 1.05)

    # 목표 섹션
    _add_text_box(slide, style.x(0.06), style.y(0.173), style.x(0.9), style.y(0.067),
                  "주간 목표",
                  font_size=style.font("subhead")["size_pt"],
                  bold=True, color=style.rgb("accent_dark"),
                  font_family=style.font_family)

    card_y = style.y(0.253)
    card_w = style.x(0.285)
    _add_kpi_card(slide, style.x(0.06), card_y, card_w,
                  "매출 목표", f"{target:,}원", style)
    _add_kpi_card(slide, style.x(0.375), card_y, card_w,
                  "성장률 목표", "5%", style)
    _add_kpi_card(slide, style.x(0.69), card_y, card_w,
                  "신규 고객 목표", "+10명", style)

    # 일정
    _add_text_box(slide, style.x(0.06), style.y(0.56), style.x(0.9), style.y(0.067),
                  "주요 일정",
                  font_size=style.font("subhead")["size_pt"],
                  bold=True, color=style.rgb("accent_dark"),
                  font_family=style.font_family)

    schedule = [
        ("월", "주간 킥오프 미팅 & KPI 리뷰"),
        ("화~수", "프로모션 캠페인 기획 및 실행"),
        ("목", "중간 점검 및 전략 조정"),
        ("금", "주간 마감 및 보고서 작성"),
    ]

    row_y = style.y(0.64)
    row_h = style.y(0.06)
    for day, task in schedule:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(style.x(0.06)), Inches(row_y),
            Inches(style.x(0.877)), Inches(row_h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = style.rgb("card_fill")
        bg.line.color.rgb = style.rgb("border")
        bg.line.width = Pt(1)

        _add_text_box(slide, style.x(0.075), row_y + style.y(0.007), style.x(0.112), row_h * 0.8,
                      day, font_size=13, bold=True, color=style.rgb("accent"),
                      font_family=style.font_family)
        _add_text_box(slide, style.x(0.21), row_y + style.y(0.007), style.x(0.675), row_h * 0.8,
                      task, font_size=13, color=style.rgb("body_text"),
                      font_family=style.font_family)
        row_y += row_h + style.y(0.007)


# ──────────────────────────────────────────────
# 메인 생성 함수
# ──────────────────────────────────────────────
def create_presentation(
    analysis: dict,
    output_path: str,
    template_path: str = None,
    style: StyleConfig = None,
) -> str:
    """분석 결과를 기반으로 PPT를 생성합니다.

    Args:
        analysis: analyze_data의 분석 결과
        output_path: 출력 파일 경로
        template_path: 회사 PPT 템플릿 경로 (선택)
        style: StyleConfig (None이면 기본 스타일 사용)
    """
    if style is None:
        style = default_style()

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        _remove_existing_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(style.slide_width)
        prs.slide_height = Inches(style.slide_height)

    print("  [슬라이드 1/5] 표지 생성...")
    create_cover_slide(prs, analysis, style)

    print("  [슬라이드 2/5] 핵심 KPI 요약...")
    create_summary_slide(prs, analysis, style)

    print("  [슬라이드 3/5] 상세 분석 차트...")
    create_chart_slide(prs, analysis, style)

    print("  [슬라이드 4/5] 주요 이슈 & 액션...")
    create_issues_slide(prs, analysis, style)

    print("  [슬라이드 5/5] 다음 주 계획...")
    create_plan_slide(prs, analysis, style)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"\n  PPT 저장 완료: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="PPT 보고서 생성")
    parser.add_argument("--analysis", required=True, help="analysis.json 경로")
    parser.add_argument("--charts", default="output/charts", help="차트 디렉토리")
    parser.add_argument("--template", default=None, help="브랜드 템플릿 .pptx 경로")
    parser.add_argument("--output", default="output/weekly_report.pptx", help="출력 파일 경로")
    args = parser.parse_args()

    with open(args.analysis, "r", encoding="utf-8") as f:
        analysis = json.load(f)

    style = default_style()
    if args.template:
        from template_analyzer import analyze_template
        from style_config import from_template_info
        info = analyze_template(args.template)
        style = from_template_info(info)

    print("[PPT 생성 시작]")
    create_presentation(analysis, args.output, args.template, style)


if __name__ == "__main__":
    main()
